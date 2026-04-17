from fastapi import FastAPI, Form
from fastapi.responses import HTMLResponse, FileResponse
from pptx import Presentation
import google.generativeai as genai
import os

app = FastAPI()

# Configuração da API do Gemini (Substitua pela sua chave real)
CHAVE_API_GEMINI = "AIzaSyDzpFWchLxj-Vkq9hjmspYhj6NZUF9BnEY"
genai.configure(api_key=CHAVE_API_GEMINI)


# 1. Front-end do MVP
@app.get("/", response_class=HTMLResponse)
async def get_form():
    return """
    <html>
        <head>
            <title>Gerador de Apresentações com IA - MVP</title>
            <style>
                body { font-family: Arial, sans-serif; padding: 50px; background-color: #f4f4f9; }
                .container { background: white; padding: 30px; border-radius: 8px; box-shadow: 0 4px 8px rgba(0,0,0,0.1); max-width: 600px; margin: auto; }
                textarea { width: 100%; padding: 10px; margin-bottom: 20px; border: 1px solid #ccc; border-radius: 4px; font-family: monospace; }
                button { padding: 10px 20px; background: #6200ea; color: white; border: none; border-radius: 4px; cursor: pointer; font-size: 16px; font-weight: bold; }
                button:hover { background: #3700b3; }
                .tag-ia { display: inline-block; background: #e0f7fa; color: #00838f; padding: 4px 8px; border-radius: 12px; font-size: 12px; margin-bottom: 15px; font-weight: bold; }
            </style>
        </head>
        <body>
            <div class="container">
                <h2>Gerador de Sprint Review 🚀</h2>
                <span class="tag-ia">✨ Powered by Gemini AI</span>
                <p>Cole logs de commits confusos ou anotações brutas. A IA irá interpretar, resumir e gerar o slide perfeitamente formatado.</p>
                <form action="/gerar-pptx" method="post">
                    <textarea name="texto_bruto" rows="8" placeholder="Ex: \nfix: resolve erro 500 no bd\nupdate: ajusta cor do botao na UI\nfeat: integra servico de email..."></textarea><br>
                    <button type="submit">Processar com IA e Gerar PPTX</button>
                </form>
            </div>
        </body>
    </html>
    """


# 2. Back-end: Processamento LLM e Geração do PPTX
@app.post("/gerar-pptx")
async def gerar_pptx(texto_bruto: str = Form(...)):
    # --- ETAPA 1: PROCESSAMENTO COM IA ---
    # Inicializa o modelo (Flash é o mais rápido e barato para tarefas de texto curtas)
    model = genai.GenerativeModel('gemini-1.5-flash')

    # Engenharia de Prompt para garantir que a IA retorne o formato exato que precisamos
    prompt = f"""
    Você é um Tech Lead analisando anotações brutas de uma equipe de desenvolvimento.
    Resuma as seguintes atividades em 3 a 5 tópicos profissionais e curtos para serem apresentados em um slide de Sprint Review para stakeholders.
    Regras estritas: 
    - Retorne APENAS os tópicos, um por linha.
    - NÃO use asteriscos, números, hífens ou marcadores no início da frase.
    - Vá direto ao ponto.

    Texto bruto:
    {texto_bruto}
    """

    # Chamada à API
    resposta_ia = model.generate_content(prompt)

    # Pega o texto gerado e divide em uma lista de tópicos
    topicos_processados = resposta_ia.text.strip().split('\n')

    # --- ETAPA 2: GERAÇÃO DO ARQUIVO PPTX ---
    prs = Presentation()

    # Capa
    slide_capa = prs.slides.add_slide(prs.slide_layouts[0])
    slide_capa.shapes.title.text = "Sprint Review"
    slide_capa.placeholders[1].text = "Resumo Inteligente gerado via IA"

    # Slide de Conteúdo
    slide_conteudo = prs.slides.add_slide(prs.slide_layouts[1])
    slide_conteudo.shapes.title.text = "Principais Entregas"
    corpo_texto = slide_conteudo.placeholders[1].text_frame

    # A primeira linha precisa ser inserida no parágrafo padrão
    p0 = corpo_texto.paragraphs[0]
    p0.text = "Destaques da iteração:"

    # Adiciona os tópicos gerados pela IA como bullet points
    for topico in topicos_processados:
        texto_limpo = topico.strip()
        if texto_limpo:
            p = corpo_texto.add_paragraph()
            p.text = texto_limpo
            p.level = 1

    file_path = "Sprint_Review_Inteligente.pptx"
    prs.save(file_path)

    return FileResponse(
        path=file_path,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        filename="Sprint_Review_Inteligente.pptx"
    )